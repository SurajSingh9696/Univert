#!/usr/bin/env python3
"""
PDF to PowerPoint converter.
Converts each PDF page to an image and creates a PPTX with images as slides.
Usage: python pdf_to_ppt.py <input_pdf> <output_pptx>
"""

import sys
import os
import tempfile
from pathlib import Path

def convert_pdf_to_ppt(pdf_path, pptx_path):
    """Convert PDF to PowerPoint by converting pages to images."""
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
        
        # Convert PDF pages to images
        print(f"Converting PDF pages to images...")
        images = convert_from_path(pdf_path, dpi=150)
        
        if not images:
            print("No pages found in PDF", file=sys.stderr)
            return False
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add blank slide layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        
        with tempfile.TemporaryDirectory() as tmpdir:
            for i, image in enumerate(images):
                print(f"Processing page {i+1}/{len(images)}...")
                
                # Save image temporarily
                img_path = os.path.join(tmpdir, f"page_{i}.png")
                image.save(img_path, "PNG")
                
                # Add slide with image
                slide = prs.slides.add_slide(blank_layout)
                
                # Add image to fill the slide
                slide.shapes.add_picture(
                    img_path,
                    Inches(0),
                    Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )
        
        # Save presentation
        prs.save(pptx_path)
        print(f"Successfully converted {pdf_path} to {pptx_path} ({len(images)} slides)")
        return True
        
    except ImportError as e:
        print(f"Missing required library: {e}", file=sys.stderr)
        return False
    except Exception as e:
        print(f"Error converting PDF to PPT: {e}", file=sys.stderr)
        return False

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_ppt.py <input_pdf> <output_pptx>", file=sys.stderr)
        sys.exit(1)
    
    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]
    
    success = convert_pdf_to_ppt(input_pdf, output_pptx)
    sys.exit(0 if success else 1)
